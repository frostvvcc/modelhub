"""
rag/document_parser 文档解析单元测试
使用临时文件测试各格式解析
"""
import os
import tempfile
import pytest
from app.services.rag.document_parser import extract_text, DocumentParser


def _openpyxl_write_broken() -> bool:
    """检测 openpyxl Workbook 写入是否受 numpy.bool 移除问题影响。"""
    try:
        import openpyxl
        import io
        wb = openpyxl.Workbook()
        wb.save(io.BytesIO())
        return False
    except Exception:
        return True


_skip_openpyxl_write = pytest.mark.skipif(
    _openpyxl_write_broken(),
    reason="openpyxl Workbook write broken (numpy.bool removed in NumPy 2.x)"
)


class TestExtractText:
    def test_nonexistent_file_returns_none(self):
        result = extract_text("/nonexistent/path/file.txt")
        assert result is None

    def test_plain_text_file(self):
        content = "这是测试文本\n第二行\n第三行"
        with tempfile.NamedTemporaryFile(mode='w', suffix='.txt', delete=False,
                                         encoding='utf-8') as f:
            f.write(content)
            tmp_path = f.name
        try:
            result = extract_text(tmp_path)
            assert result is not None
            assert "测试文本" in result
        finally:
            os.unlink(tmp_path)

    def test_markdown_file(self):
        content = "# 标题\n\n这是**内容**\n\n## 子标题\n\n段落文字"
        with tempfile.NamedTemporaryFile(mode='w', suffix='.md', delete=False,
                                         encoding='utf-8') as f:
            f.write(content)
            tmp_path = f.name
        try:
            result = extract_text(tmp_path)
            assert result is not None
            assert "标题" in result
        finally:
            os.unlink(tmp_path)

    def test_markdown_removes_html_comments(self):
        content = "有效内容\n<!-- 这是注释 -->\n更多内容"
        with tempfile.NamedTemporaryFile(mode='w', suffix='.md', delete=False,
                                         encoding='utf-8') as f:
            f.write(content)
            tmp_path = f.name
        try:
            result = extract_text(tmp_path)
            assert result is not None
            assert "这是注释" not in result
            assert "有效内容" in result
        finally:
            os.unlink(tmp_path)

    @_skip_openpyxl_write
    def test_excel_file(self):
        pytest.importorskip("openpyxl")
        import openpyxl
        wb = openpyxl.Workbook()
        ws = wb.active
        ws.title = "测试表"
        ws.append(["姓名", "成绩", "备注"])
        ws.append(["张三", 90, "优秀"])
        ws.append(["李四", 75, "良好"])

        with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as f:
            tmp_path = f.name
        wb.save(tmp_path)
        try:
            result = extract_text(tmp_path)
            assert result is not None
            assert "测试表" in result
            assert "张三" in result
            assert "成绩" in result
        finally:
            os.unlink(tmp_path)

    @_skip_openpyxl_write
    def test_excel_multiple_sheets(self):
        pytest.importorskip("openpyxl")
        import openpyxl
        wb = openpyxl.Workbook()
        ws1 = wb.active
        ws1.title = "Sheet1"
        ws1.append(["数据A"])
        ws2 = wb.create_sheet("Sheet2")
        ws2.append(["数据B"])

        with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as f:
            tmp_path = f.name
        wb.save(tmp_path)
        try:
            result = extract_text(tmp_path)
            assert result is not None
            assert "Sheet1" in result
            assert "Sheet2" in result
            assert "数据A" in result
            assert "数据B" in result
        finally:
            os.unlink(tmp_path)

    def test_pptx_file(self):
        pytest.importorskip("pptx")
        from pptx import Presentation
        from pptx.util import Inches
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "测试标题"
        slide.placeholders[1].text = "测试内容"

        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as f:
            tmp_path = f.name
        prs.save(tmp_path)
        try:
            result = extract_text(tmp_path)
            assert result is not None
            assert "测试标题" in result
        finally:
            os.unlink(tmp_path)

    def test_pptx_with_notes(self):
        pytest.importorskip("pptx")
        from pptx import Presentation
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        # 添加文本框
        from pptx.util import Inches, Pt
        txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(2))
        txBox.text_frame.text = "幻灯片正文"
        # 添加备注
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = "这是备注内容"

        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as f:
            tmp_path = f.name
        prs.save(tmp_path)
        try:
            result = extract_text(tmp_path)
            assert result is not None
            assert "备注" in result
        finally:
            os.unlink(tmp_path)


class TestDocumentParserClass:
    def test_extract_via_class(self):
        content = "通过 DocumentParser 类解析"
        with tempfile.NamedTemporaryFile(mode='w', suffix='.txt', delete=False,
                                         encoding='utf-8') as f:
            f.write(content)
            tmp_path = f.name
        try:
            result = DocumentParser.extract(tmp_path)
            assert result is not None
            assert "DocumentParser" in result
        finally:
            os.unlink(tmp_path)


# ------------------------------------------------------------------
# Word 表格结构化提取
# ------------------------------------------------------------------

class TestDocxTableExtraction:
    def _build_docx(self, tmp_path):
        from docx import Document
        doc = Document()
        doc.add_paragraph("表格前的说明文字")
        table = doc.add_table(rows=2, cols=2)
        table.rows[0].cells[0].text = "课程"
        table.rows[0].cells[1].text = "学分"
        table.rows[1].cells[0].text = "操作系统"
        table.rows[1].cells[1].text = "4"
        doc.add_paragraph("表格后的结论文字")
        path = tmp_path / "sample.docx"
        doc.save(str(path))
        return str(path)

    def test_table_preserved_as_markdown_in_position(self, tmp_path):
        from app.services.rag.document_parser import _extract_docx
        text = _extract_docx(self._build_docx(tmp_path))
        assert text is not None
        assert "| 课程 | 学分 |" in text
        assert "| --- | --- |" in text
        assert "| 操作系统 | 4 |" in text
        # 表格位于前后两段文字之间，位置保留
        assert text.index("表格前的说明文字") < text.index("| 课程 |") < text.index("表格后的结论文字")
