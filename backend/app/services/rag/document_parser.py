"""
文档解析模块
支持 PDF、Word（.docx/.doc）、纯文本、Excel、PPT、Markdown、图片（OCR）
包含网页爬取文本的 Boilerplate Removal（正文提取 / 噪声清洗）
"""
import os
import re
import tempfile
from typing import Optional, List
from app.utils.logger_config import get_logger

logger = get_logger(__name__)


# ------------------------------------------------------------------
# 文本清洗（Boilerplate Removal）
# ------------------------------------------------------------------

_METADATA_PATTERN = re.compile(
    r'^(标题|来源URL|发布日期|抓取时间|来源|发布时间)\s*[:：].*$',
    re.MULTILINE,
)

_FOOTER_KEYWORDS = frozenset({
    "icp备", "all rights reserved", "版权所有", "备案号", "鲁公网安备",
    "邮编", "copyright", "鲁icp",
})

_NAV_KEYWORDS = frozenset({
    "首页", "学校概况", "学校简介", "大学章程", "历史沿革", "学校文化",
    "现任领导", "历任领导", "组织机构", "人才培养", "本科生教育", "研究生教育",
    "国际生教育", "远程教育", "学科建设", "招生就业", "本科生招生", "研究生招生",
    "成人招生", "就业网", "科学研究", "科研工作", "人文社科", "技术转移",
    "人才招聘", "合作交流", "校友之家", "基金会", "校友办", "新闻网",
    "教学单位", "机构设置", "导航", "正文", "信息公开", "接诉即办",
    "党政办公室", "教务部", "学团工作部", "安全管理部", "后勤管理部",
    "图书馆", "国际交流部", "工会", "科研部", "财务部", "资产部",
})

_ATTACHMENT_PATTERN = re.compile(r'附件【[^】]*】已下载\s*\d*\s*次?')


def clean_web_text(text: str) -> str:
    """
    清洗网页爬取文本，去除非正文内容。
    处理 5 种噪声：元数据头、导航菜单、重复导航、页脚版权、附件标注。
    """
    if not text:
        return ""

    # 1. 去元数据头行（标题:/来源URL:/发布日期:/抓取时间:）
    text = _METADATA_PATTERN.sub('', text)

    # 2. 去附件下载标注
    text = _ATTACHMENT_PATTERN.sub('', text)

    # 3. 逐行过滤：去导航菜单行和页脚行
    lines = text.split('\n')
    cleaned_lines: List[str] = []
    nav_streak = 0

    for line in lines:
        stripped = line.strip()
        if not stripped:
            if cleaned_lines and cleaned_lines[-1] != '':
                cleaned_lines.append('')
            nav_streak = 0
            continue

        lower = stripped.lower()

        # 页脚版权行
        if any(kw in lower for kw in _FOOTER_KEYWORDS):
            continue

        # 导航菜单行：短行（<=10字符）且是导航关键词
        if len(stripped) <= 10 and stripped in _NAV_KEYWORDS:
            nav_streak += 1
            continue

        # 如果之前连续跳过了多个导航行，说明在菜单区域
        if nav_streak >= 3 and len(stripped) <= 10:
            nav_streak += 1
            continue

        nav_streak = 0

        # 单个字符的噪声行（如 "X", "-"）
        if len(stripped) <= 1:
            continue

        cleaned_lines.append(stripped)

    # 4. 去重复段落（相同内容只保留第一次出现）
    seen = set()
    deduped: List[str] = []
    for line in cleaned_lines:
        if line == '' or line not in seen:
            deduped.append(line)
            if line:
                seen.add(line)

    result = '\n'.join(deduped).strip()

    # 5. 合并连续空行
    result = re.sub(r'\n{3,}', '\n\n', result)

    if text and result:
        ratio = len(result) / len(text)
        logger.info(f"文本清洗: 原始 {len(text)} 字符 → 清洗后 {len(result)} 字符 (保留 {ratio:.0%})")

    return result


def is_worth_indexing(text: str, min_length: int = 100) -> bool:
    """判断清洗后的文本是否值得入库"""
    if not text or len(text.strip()) < min_length:
        return False
    meaningful_chars = sum(1 for c in text if c.isalnum() or c in '，。！？、；：""''（）')
    return meaningful_chars >= min_length * 0.5


def extract_text(file_path: str) -> Optional[str]:
    """
    统一入口：根据文件扩展名分发到对应的解析器。

    Args:
        file_path: 文件路径

    Returns:
        提取的文本内容，失败时返回 None
    """
    if not os.path.exists(file_path):
        logger.error(f"文件不存在: {file_path}")
        return None

    ext = os.path.splitext(file_path)[1].lower()

    if ext in ('.docx', '.doc'):
        return _extract_docx(file_path)
    elif ext == '.pdf':
        return _extract_pdf(file_path)
    elif ext in ('.xlsx', '.xls'):
        return _extract_excel(file_path)
    elif ext in ('.pptx', '.ppt'):
        return _extract_pptx(file_path)
    elif ext == '.md':
        return _extract_markdown(file_path)
    else:
        # 纯文本（.txt、.csv 等）
        return _extract_plain_text(file_path)


# ------------------------------------------------------------------
# Word (.docx / .doc)
# ------------------------------------------------------------------

def _extract_docx(file_path: str) -> Optional[str]:
    """提取 Word 文档文本（支持 .docx 和 .doc）"""
    ext = os.path.splitext(file_path)[1].lower()
    try:
        if ext == '.doc':
            # 优先用 doc2text
            try:
                import doc2text
                text = doc2text.extract(file_path)
                if text and text.strip():
                    return text.strip()
            except (ImportError, Exception) as e:
                logger.warning(f"doc2text 处理失败（{e}），尝试 python-docx")

        # .docx 或 doc2text 失败后的 fallback
        from docx import Document
        doc = Document(file_path)
        paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
        for table in doc.tables:
            for row in table.rows:
                row_text = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if row_text:
                    paragraphs.append(' | '.join(row_text))
        text = '\n'.join(paragraphs)
        logger.info(f"Word 文档解析完成: {len(paragraphs)} 段落, {len(text)} 字符")
        return text or None
    except Exception as e:
        logger.error(f"Word 文档解析失败 ({file_path}): {e}")
        return None


# ------------------------------------------------------------------
# PDF
# ------------------------------------------------------------------

def _extract_pdf_with_marker(file_path: str) -> Optional[str]:
    """用 Marker 将 PDF 高保真转为结构化 Markdown（保留表格、标题层级、公式）"""
    try:
        from marker.converters.pdf import PdfConverter
        # marker >= 1.0: PdfConverter() 自动加载模型
        converter = PdfConverter()
        rendered = converter(file_path)

        # marker >= 1.0 返回 Document 对象，有 .markdown 属性
        # marker 0.x 返回 tuple (text, images, metadata)
        if hasattr(rendered, 'markdown'):
            md_text = rendered.markdown.strip()
        elif isinstance(rendered, tuple) and len(rendered) >= 1:
            md_text = str(rendered[0]).strip()
        else:
            md_text = ""

        if md_text and len(md_text) >= 20:
            logger.info(f"PDF 解析完成 (Marker): {len(md_text)} 字符")
            return md_text
        return None
    except ImportError:
        return None
    except Exception as e:
        logger.warning(f"Marker PDF 解析失败，降级到 pdfplumber: {e}")
        return None


def _extract_pdf_with_tables(file_path: str) -> Optional[str]:
    """用 pdfplumber 提取文本 + 表格（表格转 Markdown 格式，排除表格区域避免重复）"""
    try:
        import pdfplumber
    except ImportError:
        return None

    parts = []
    with pdfplumber.open(file_path) as pdf:
        for page_idx, page in enumerate(pdf.pages):
            tables = page.find_tables() or []
            table_bboxes = [t.bbox for t in tables]

            # 提取表格并转 Markdown
            for table in tables:
                table_data = table.extract()
                if not table_data:
                    continue
                rows = []
                for ri, row in enumerate(table_data):
                    cells = [str(c or '').strip() for c in row]
                    rows.append('| ' + ' | '.join(cells) + ' |')
                    if ri == 0:
                        rows.append('| ' + ' | '.join(['---'] * len(cells)) + ' |')
                parts.append('\n'.join(rows))

            # 提取非表格区域的文本，避免内容重复
            if table_bboxes:
                filtered_page = page
                for bbox in table_bboxes:
                    filtered_page = filtered_page.outside_bbox(bbox)
                text = filtered_page.extract_text()
            else:
                text = page.extract_text()

            if text and text.strip():
                parts.append(text.strip())

    result = '\n\n'.join(parts).strip()
    if result:
        logger.info(f"PDF 解析完成 (pdfplumber+tables): {len(result)} 字符")
    return result or None


def _extract_pdf(file_path: str) -> Optional[str]:
    """提取 PDF 文本：Marker（高保真 Markdown）→ pdfplumber+tables → PyPDF2 → OCR"""
    try:
        # 优先：Marker 高保真 PDF→Markdown（保留表格/标题/公式）
        marker_text = _extract_pdf_with_marker(file_path)
        if marker_text:
            return marker_text

        # 次选：pdfplumber 带表格提取
        plumber_text = _extract_pdf_with_tables(file_path)
        if plumber_text and len(plumber_text) >= 20:
            return plumber_text

        # 降级：PyPDF2 纯文本
        text_from_layer = None
        try:
            import PyPDF2
            parts = []
            with open(file_path, 'rb') as f:
                reader = PyPDF2.PdfReader(f)
                for page in reader.pages:
                    text = page.extract_text()
                    if text:
                        parts.append(text)
            text_from_layer = '\n'.join(parts).strip()
            if text_from_layer:
                logger.info(f"PDF 解析完成 (PyPDF2): {len(text_from_layer)} 字符")
        except ImportError:
            pass

        if text_from_layer and len(text_from_layer) >= 20:
            return text_from_layer

        # 最后兜底：OCR
        ocr_text = _extract_pdf_with_ocr(file_path)
        if ocr_text:
            return ocr_text
        return text_from_layer or None
    except Exception as e:
        logger.error(f"PDF 解析失败 ({file_path}): {e}")
        return None


def _extract_pdf_with_ocr(file_path: str) -> Optional[str]:
    """将扫描版 PDF 页面渲染成图片后交给 PaddleOCR 识别。"""
    try:
        import fitz  # PyMuPDF
    except ImportError:
        logger.warning("PyMuPDF 未安装，无法对扫描版 PDF 做 OCR。可运行: pip install PyMuPDF")
        return None

    try:
        from app.utils.ocr_utils import PADDLEOCR_AVAILABLE, get_ocr_processor
        if not PADDLEOCR_AVAILABLE:
            logger.warning("PaddleOCR 未安装，扫描版 PDF OCR 不可用")
            return None

        processor = get_ocr_processor()
        if not processor:
            logger.warning("PaddleOCR 初始化失败，扫描版 PDF OCR 不可用")
            return None

        parts = []
        with tempfile.TemporaryDirectory(prefix="pdf_ocr_") as temp_dir:
            with fitz.open(file_path) as doc:
                page_count = len(doc)
                for page_index in range(page_count):
                    page = doc.load_page(page_index)
                    pixmap = page.get_pixmap(matrix=fitz.Matrix(2, 2), alpha=False)
                    image_path = os.path.join(temp_dir, f"page_{page_index + 1}.png")
                    pixmap.save(image_path)
                    page_text = processor.extract_text_from_image(image_path)
                    if page_text:
                        parts.append(f"第 {page_index + 1} 页\n{page_text.strip()}")

        full_text = "\n\n".join(parts).strip()
        logger.info(f"PDF OCR 解析完成: {len(full_text)} 字符")
        return full_text or None
    except Exception as e:
        logger.error(f"PDF OCR 解析失败 ({file_path}): {e}")
        return None


# ------------------------------------------------------------------
# Excel (.xlsx / .xls)
# ------------------------------------------------------------------

def _extract_excel(file_path: str) -> Optional[str]:
    """
    提取 Excel 文本。
    每个 Sheet 输出为「Sheet名：\n表头 | 列1 | 列2 ...」格式，保留表头语义。
    """
    ext = os.path.splitext(file_path)[1].lower()
    try:
        if ext == '.xlsx':
            import openpyxl
            wb = openpyxl.load_workbook(file_path, data_only=True)
            parts = []
            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                rows = list(ws.iter_rows(values_only=True))
                if not rows:
                    continue
                parts.append(f"【{sheet_name}】")
                for row in rows:
                    row_text = [str(c) if c is not None else '' for c in row]
                    if any(t.strip() for t in row_text):
                        parts.append(' | '.join(row_text))
            text = '\n'.join(parts)
            logger.info(f"Excel 解析完成: {len(wb.sheetnames)} 个 Sheet, {len(text)} 字符")
            return text or None
        else:
            # .xls
            import xlrd
            wb = xlrd.open_workbook(file_path)
            parts = []
            for sheet in wb.sheets():
                parts.append(f"【{sheet.name}】")
                for r in range(sheet.nrows):
                    row_text = [str(sheet.cell_value(r, c)) for c in range(sheet.ncols)]
                    if any(t.strip() for t in row_text):
                        parts.append(' | '.join(row_text))
            text = '\n'.join(parts)
            logger.info(f"XLS 解析完成: {len(text)} 字符")
            return text or None
    except ImportError as e:
        logger.error(f"Excel 解析库未安装: {e}，请运行 pip install openpyxl xlrd")
        return None
    except Exception as e:
        logger.error(f"Excel 解析失败 ({file_path}): {e}")
        return None


# ------------------------------------------------------------------
# PowerPoint (.pptx / .ppt)
# ------------------------------------------------------------------

def _extract_pptx(file_path: str) -> Optional[str]:
    """
    提取 PPT 文本。
    每张幻灯片提取文本框内容和备注，图片部分走 OCR（如可用）。
    """
    ext = os.path.splitext(file_path)[1].lower()
    if ext == '.ppt':
        logger.warning(f"旧版 .ppt 格式不支持直接解析，请转换为 .pptx")
        return None
    try:
        from pptx import Presentation
        from pptx.util import Inches
        prs = Presentation(file_path)
        parts = []
        for slide_idx, slide in enumerate(prs.slides, 1):
            slide_texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            slide_texts.append(text)
            # 备注
            if slide.has_notes_slide:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    slide_texts.append(f"[备注] {notes}")
            if slide_texts:
                parts.append(f"【第{slide_idx}页】")
                parts.extend(slide_texts)
        text = '\n'.join(parts)
        logger.info(f"PPT 解析完成: {len(prs.slides)} 页, {len(text)} 字符")
        return text or None
    except ImportError:
        logger.error("python-pptx 未安装，请运行 pip install python-pptx")
        return None
    except Exception as e:
        logger.error(f"PPT 解析失败 ({file_path}): {e}")
        return None


# ------------------------------------------------------------------
# Markdown (.md)
# ------------------------------------------------------------------

def _extract_markdown(file_path: str) -> Optional[str]:
    """
    解析 Markdown 文件。
    保留标题层级结构（作为后续分块策略的边界标记），去除 HTML 标签。
    """
    try:
        from app.utils.encoding_utils import read_file_smart
        text = read_file_smart(file_path)
        if not text:
            return None
        # 去除 HTML 注释
        import re
        text = re.sub(r'<!--.*?-->', '', text, flags=re.DOTALL)
        logger.info(f"Markdown 解析完成: {len(text)} 字符")
        return text.strip() or None
    except Exception as e:
        logger.error(f"Markdown 解析失败 ({file_path}): {e}")
        return None


# ------------------------------------------------------------------
# 纯文本
# ------------------------------------------------------------------

def _extract_plain_text(file_path: str) -> Optional[str]:
    """读取纯文本文件（自动检测编码）"""
    try:
        from app.utils.encoding_utils import read_file_smart
        return read_file_smart(file_path)
    except Exception as e:
        logger.error(f"纯文本读取失败 ({file_path}): {e}")
        return None


# ------------------------------------------------------------------
# 向后兼容类（供 vector_service.py 调用）
# ------------------------------------------------------------------

class DocumentParser:
    """文档解析器类（封装，方便将来扩展）"""

    @staticmethod
    def extract(file_path: str) -> Optional[str]:
        return extract_text(file_path)
